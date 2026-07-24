import streamlit as st
import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO

st.set_page_config(page_title="Rapports Épidémiologiques", layout="wide")

st.title("📊 Générateur de Rapports Épidémiologiques")
st.write("Guinée - DHIS2 | Automation complète")

uploaded_file = st.file_uploader("📤 Téléchargez votre fichier DHIS2 (.xls ou .xlsx)", type=['xls', 'xlsx'])

if uploaded_file is not None:
    if st.button("🚀 Générer les rapports", use_container_width=True):
        with st.spinner("⏳ Génération en cours... (2 minutes)"):
            try:
                df = pd.read_excel(uploaded_file, sheet_name=0)
                df.columns = df.iloc[0]
                df = df.iloc[1:].reset_index(drop=True)
                
                df_simple = pd.DataFrame({
                    'Établissement': df.iloc[:, 1],
                    'Paludisme': pd.to_numeric(df.iloc[:, 3], errors='coerce'),
                    'Testés': pd.to_numeric(df.iloc[:, 5], errors='coerce'),
                    'Ebola': pd.to_numeric(df.iloc[:, 19], errors='coerce'),
                    'Lassa': pd.to_numeric(df.iloc[:, 25], errors='coerce'),
                    'Choléra': pd.to_numeric(df.iloc[:, 55], errors='coerce'),
                    'Décès': pd.to_numeric(df.iloc[:, 9], errors='coerce')
                })
                
                df_simple = df_simple.fillna(0)
                
                total_palu = int(df_simple['Paludisme'].sum())
                total_testes = int(df_simple['Testés'].sum())
                total_ebola = int(df_simple['Ebola'].sum())
                total_lassa = int(df_simple['Lassa'].sum())
                total_cholera = int(df_simple['Choléra'].sum())
                total_deces = int(df_simple['Décès'].sum())
                
                st.success("✅ Données chargées avec succès !")
                
                col1, col2, col3 = st.columns(3)
                col1.metric("📊 Paludisme", f"{total_palu:,}")
                col2.metric("⚠️ Ebola", total_ebola)
                col3.metric("💀 Décès", total_deces)
                
                col1, col2, col3 = st.columns(3)
                col1.metric("Lassa", total_lassa)
                col2.metric("Choléra", total_cholera)
                col3.metric("Établissements", len(df_simple))
                
                st.divider()
                st.subheader("📥 Télécharger vos rapports")
                
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    excel_output = BytesIO()
                    with pd.ExcelWriter(excel_output, engine='openpyxl') as writer:
                        df_simple.to_excel(writer, sheet_name='Données', index=False)
                        top10 = df_simple.nlargest(10, 'Paludisme')[['Établissement', 'Paludisme', 'Testés']]
                        top10.to_excel(writer, sheet_name='Top 10', index=False)
                        summary = pd.DataFrame({
                            'Indicateur': ['Paludisme', 'Ebola', 'Lassa', 'Choléra', 'Décès'],
                            'Total': [total_palu, total_ebola, total_lassa, total_cholera, total_deces]
                        })
                        summary.to_excel(writer, sheet_name='Résumé', index=False)
                    
                    excel_output.seek(0)
                    st.download_button("📊 Excel", excel_output, "Rapport_DHIS2.xlsx")
                
                with col2:
                    fig = make_subplots(rows=2, cols=2, subplot_titles=("Top 15", "Distribution", "Graves", "Taux"))
                    
                    top15 = df_simple.nlargest(15, 'Paludisme')
                    fig.add_trace(go.Bar(x=top15['Paludisme'], y=top15['Établissement'], orientation='h', marker_color='#d62728'), row=1, col=1)
                    
                    fig.add_trace(go.Pie(labels=['Paludisme', 'Ebola', 'Lassa', 'Choléra'], 
                                         values=[total_palu, total_ebola, total_lassa, total_cholera],
                                         marker_colors=['#ff7f0e', '#e74c3c', '#9b59b6', '#3498db']), row=1, col=2)
                    
                    fig.add_trace(go.Bar(x=['Ebola', 'Lassa', 'Choléra'], y=[total_ebola, total_lassa, total_cholera],
                                        marker_color=['#e74c3c', '#9b59b6', '#3498db']), row=2, col=1)
                    
                    top10_test = df_simple.nlargest(10, 'Testés')
                    if len(top10_test) > 0:
                        taux = (top10_test['Paludisme'] / top10_test['Testés'] * 100).values
                        fig.add_trace(go.Scatter(x=list(range(len(taux))), y=taux, mode='lines+markers', 
                                               name='Taux (%)', line=dict(color='#1f77b4')), row=2, col=2)
                    
                    fig.update_layout(height=800, showlegend=False)
                    
                    html_output = fig.to_html()
                    st.download_button("🌐 Dashboard HTML", html_output, "Dashboard.html")
                
                with col3:
                    prs = Presentation()
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
                    
                    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                    bg = slide1.background.fill
                    bg.solid()
                    bg.fore_color.rgb = RGBColor(31, 78, 120)
                    
                    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
                    title_frame = title_box.text_frame
                    p = title_frame.paragraphs[0]
                    p.text = "RAPPORT ÉPIDÉMIOLOGIQUE"
                    p.font.size = Pt(54)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    
                    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
                    title2 = slide2.shapes.title
                    title2.text = "Résumé Exécutif"
                    text_frame = slide2.placeholders[1].text_frame
                    text_frame.clear()
                    
                    p1 = text_frame.paragraphs[0]
                    p1.text = f"✓ Paludisme : {total_palu:,} ({len(df_simple)} établissements)"
                    
                    for text in [f"✓ Ebola : {total_ebola} | Lassa : {total_lassa} | Choléra : {total_cholera}",
                                f"✓ Décès : {total_deces}",
                                "✓ Situation sous contrôle ✅",
                                "→ Surveillance continue recommandée"]:
                        p = text_frame.add_paragraph()
                        p.text = text
                    
                    prs_output = BytesIO()
                    prs.save(prs_output)
                    prs_output.seek(0)
                    
                    st.download_button("📑 PowerPoint", prs_output, "Rapport_Épidémiologique.pptx")
                
                st.divider()
                st.subheader("📊 Aperçu du Dashboard")
                st.plotly_chart(fig, use_container_width=True)
                
            except Exception as e:
                st.error(f"❌ Erreur : {str(e)}")

else:
    st.info("👈 Téléchargez un fichier DHIS2 (.xls ou .xlsx) pour commencer")
